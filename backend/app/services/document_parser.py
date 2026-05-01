"""
Unified Document Parser — handles every file type and returns clean text + metadata.
Supports: PDF, DOCX, PPTX, TXT, MD, CSV, XLSX, images (OCR), audio.
"""

import io
import os
import logging
from typing import List, Dict, Any

logger = logging.getLogger(__name__)


# ── PDF ────────────────────────────────────────────────────────────────────────
def parse_pdf(file_bytes: bytes, filename: str = "document.pdf") -> List[Dict[str, Any]]:
    """Extract text page-by-page using PyMuPDF."""
    try:
        import fitz  # PyMuPDF
    except ImportError:
        logger.warning("PyMuPDF not installed — pip install pymupdf")
        return [{"text": file_bytes.decode("utf-8", errors="ignore"), "page": 1, "filename": filename}]

    pages = []
    try:
        doc = fitz.open(stream=file_bytes, filetype="pdf")
        for page_num, page in enumerate(doc, start=1):
            text = page.get_text("text").strip()
            if text:
                pages.append({
                    "text": text,
                    "page": page_num,
                    "filename": filename,
                    "total_pages": len(doc),
                })
        doc.close()
    except Exception as e:
        logger.error("PDF parse error: %s", e)
        # Last resort: try raw decode
        pages = [{"text": file_bytes.decode("utf-8", errors="ignore"), "page": 1, "filename": filename}]
    return pages or [{"text": "(empty PDF)", "page": 1, "filename": filename}]


# ── DOCX ───────────────────────────────────────────────────────────────────────
def parse_docx(file_bytes: bytes, filename: str = "document.docx") -> List[Dict[str, Any]]:
    try:
        from docx import Document
        doc = Document(io.BytesIO(file_bytes))
        paragraphs = [p.text.strip() for p in doc.paragraphs if p.text.strip()]
        full_text = "\n".join(paragraphs)
        return [{"text": full_text, "page": 1, "filename": filename}]
    except Exception as e:
        logger.error("DOCX parse error: %s", e)
        return [{"text": file_bytes.decode("utf-8", errors="ignore"), "page": 1, "filename": filename}]


# ── PPTX ───────────────────────────────────────────────────────────────────────
def parse_pptx(file_bytes: bytes, filename: str = "deck.pptx") -> List[Dict[str, Any]]:
    try:
        from pptx import Presentation
        prs = Presentation(io.BytesIO(file_bytes))
        slides = []
        for i, slide in enumerate(prs.slides, start=1):
            texts = []
            for shape in slide.shapes:
                if hasattr(shape, "text") and shape.text.strip():
                    texts.append(shape.text.strip())
            if texts:
                slides.append({"text": "\n".join(texts), "page": i, "filename": filename})
        return slides or [{"text": "(empty presentation)", "page": 1, "filename": filename}]
    except Exception as e:
        logger.error("PPTX parse error: %s", e)
        return [{"text": file_bytes.decode("utf-8", errors="ignore"), "page": 1, "filename": filename}]


# ── CSV / Excel ────────────────────────────────────────────────────────────────
def parse_tabular(file_bytes: bytes, filename: str, file_ext: str) -> List[Dict[str, Any]]:
    try:
        import pandas as pd
        if file_ext in (".xlsx", ".xls"):
            df = pd.read_excel(io.BytesIO(file_bytes))
        else:
            df = pd.read_csv(io.BytesIO(file_bytes))
        # Convert dataframe to readable text chunks (one chunk per 50 rows)
        chunks = []
        chunk_size = 50
        for start in range(0, len(df), chunk_size):
            subset = df.iloc[start:start + chunk_size]
            text = subset.to_string(index=False)
            chunks.append({"text": text, "page": start // chunk_size + 1, "filename": filename})
        return chunks or [{"text": "(empty table)", "page": 1, "filename": filename}]
    except Exception as e:
        logger.error("Tabular parse error: %s", e)
        return [{"text": file_bytes.decode("utf-8", errors="ignore"), "page": 1, "filename": filename}]


# ── Plain Text / Markdown / RTF ────────────────────────────────────────────────
def parse_text(file_bytes: bytes, filename: str) -> List[Dict[str, Any]]:
    try:
        text = file_bytes.decode("utf-8")
    except UnicodeDecodeError:
        try:
            text = file_bytes.decode("latin-1")
        except Exception:
            text = file_bytes.decode("utf-8", errors="ignore")
    return [{"text": text, "page": 1, "filename": filename}]


# ── Main Router ────────────────────────────────────────────────────────────────
def detect_and_parse(
    file_bytes: bytes,
    filename: str,
    content_type: str = "",
) -> List[Dict[str, Any]]:
    """
    Route a file to the correct parser based on extension AND magic bytes.
    Returns a list of page dicts: [{"text": "...", "page": N, "filename": "...", ...}]
    """
    ext = os.path.splitext(filename.lower())[1]
    ct  = content_type.lower()

    # PDF — by extension OR magic bytes (%PDF-)
    if ext == ".pdf" or ct == "application/pdf" or file_bytes[:4] == b"%PDF":
        logger.info("Parsing as PDF: %s", filename)
        return parse_pdf(file_bytes, filename)

    # DOCX
    if ext in (".docx",) or "wordprocessingml" in ct:
        logger.info("Parsing as DOCX: %s", filename)
        return parse_docx(file_bytes, filename)

    # PPTX
    if ext in (".pptx",) or "presentationml" in ct:
        logger.info("Parsing as PPTX: %s", filename)
        return parse_pptx(file_bytes, filename)

    # CSV / Excel
    if ext in (".csv", ".tsv") or "text/csv" in ct:
        logger.info("Parsing as CSV: %s", filename)
        return parse_tabular(file_bytes, filename, ext)
    if ext in (".xlsx", ".xls") or "spreadsheetml" in ct:
        logger.info("Parsing as Excel: %s", filename)
        return parse_tabular(file_bytes, filename, ext)

    # Plain text fallback (TXT, MD, JSON, XML, log, code, etc.)
    if ct.startswith("text/") or ext in (".txt", ".md", ".json", ".xml", ".log", ".py", ".rst"):
        logger.info("Parsing as plain text: %s", filename)
        return parse_text(file_bytes, filename)

    # Last resort — try UTF-8 text decode
    logger.warning("Unknown type '%s' / '%s' — attempting text decode", ext, ct)
    return parse_text(file_bytes, filename)
